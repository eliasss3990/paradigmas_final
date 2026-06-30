"""Genera presentacion_G15.pptx (11 slides) para la defensa del Proyecto Final.

Estilo: demo de producto (UAT) — paleta verde/azul/ámbar, tarjetas con barra lateral.
Las "capturas" provienen de ./mockups (generadas con gen_mockups.py).

Uso:
    python gen_mockups.py    # 1) genera las capturas en ./mockups
    python build_ppt.py      # 2) genera ../docs/presentacion_G15.pptx
(requiere python-pptx y Pillow; ver README.md)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
MK = os.path.join(HERE, "mockups")
OUT = os.path.join(HERE, "..", "docs", "presentacion_G15.pptx")

# --- Datos editables (cambiar acá si varía la info) ---
INTEGRANTES = ["Elías González", "Enzo Domínguez", "Nicolás Bareiro"]
PIE_CIERRE = "Grupo 15 · " + " · ".join(INTEGRANTES) + " · FP-UNA 2026"
FOOTER = "Sistema de Gestión de Eventos · Proyecto Final · FP-UNA 2026"

# Paleta (inspirada en plantilla del compañero)
DARK   = RGBColor(0x0B, 0x12, 0x20)
EMER   = RGBColor(0x10, 0xB9, 0x81)
EMER_D = RGBColor(0x04, 0x78, 0x57)
BLUE   = RGBColor(0x1E, 0x3A, 0x8A)
BLUE2  = RGBColor(0x25, 0x63, 0xEB)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
BG     = RGBColor(0xF8, 0xFA, 0xFC)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x0B, 0x12, 0x20)
BODY   = RGBColor(0x11, 0x18, 0x27)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
LINE   = RGBColor(0xE2, 0xE8, 0xF0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SUBT   = RGBColor(0xE0, 0xF2, 0xFE)
FAINT  = RGBColor(0xD1, 0xFA, 0xE5)
FH = "Segoe UI Semibold"; FB = "Segoe UI"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(): return prs.slides.add_slide(BLANK)
def bg(s, c):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c
def rect(s, x, y, w, h, c, line=None, lw=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = c
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = lw or Pt(1)
    sp.shadow.inherit = False
    return sp
def P(t, sz, c, b=False, f=FB): return [(t, sz, c, b, f)]
def tb(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sa=Pt(4), ls=1.0):
    box = s.shapes.add_textbox(x, y, w, h); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor; tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = sa; p.space_before = Pt(0); p.line_spacing = ls
        for (t, sz, c, b, f) in para:
            r = p.add_run(); r.text = t; r.font.size = sz; r.font.color.rgb = c
            r.font.bold = b; r.font.name = f
    return box
def header(s, kicker, title, num):
    rect(s, 0, 0, Inches(0.2), Inches(7.5), EMER)
    rect(s, Inches(0.2), 0, Inches(0.08), Inches(7.5), BLUE)
    tb(s, Inches(0.6), Inches(0.32), Inches(11), Inches(0.3), [P(kicker.upper(), Pt(10), EMER_D, True)])
    tb(s, Inches(0.58), Inches(0.62), Inches(12.2), Inches(0.7), [P(title, Pt(26), INK, True, FH)])
    rect(s, Inches(0.6), Inches(1.42), Inches(12.1), Pt(1.4), LINE)
    tb(s, Inches(0.6), Inches(7.06), Inches(9), Inches(0.3), [P(FOOTER, Pt(8), MUTED)])
    tb(s, Inches(12.2), Inches(7.06), Inches(0.7), Inches(0.3), [P(f"{num:02d}", Pt(8.5), MUTED)], align=PP_ALIGN.RIGHT)
def card(s, x, y, w, h, bar, title, body, tsz=14, bsz=10.6):
    rect(s, x, y, w, h, CARD, line=LINE, lw=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, Inches(0.09), h, bar)
    tb(s, x + Inches(0.28), y + Inches(0.16), w - Inches(0.5), Inches(0.4), [P(title, Pt(tsz), INK, True, FH)])
    if isinstance(body, str): body = [body]
    paras = [P(b, Pt(bsz), BODY) for b in body]
    tb(s, x + Inches(0.28), y + Inches(0.62), w - Inches(0.5), h - Inches(0.8), paras, ls=1.05)
def pic_fit(s, path, x, y, maxw, maxh):
    iw, ih = Image.open(path).size
    w = maxw; h = Emu(int(w * ih / iw))
    if h > maxh: h = maxh; w = Emu(int(h * iw / ih))
    s.shapes.add_picture(path, x + (maxw - w)//2, y + (maxh - h)//2, width=w, height=h)


# ============ SLIDE 1 — Portada ============
s = slide(); bg(s, DARK)
rect(s, Inches(9.6), Inches(-1.0), Inches(4.5), Inches(4.5), RGBColor(0x12,0x1C,0x30), shape=MSO_SHAPE.OVAL)
rect(s, Inches(0.7), Inches(1.0), Inches(1.6), Inches(0.14), EMER)
tb(s, Inches(0.68), Inches(1.3), Inches(6.0), Inches(0.3), [P("PROYECTO FINAL · GRUPO 15", Pt(11), EMER, True)])
tb(s, Inches(0.66), Inches(1.75), Inches(7.2), Inches(2.0),
   [P("Sistema de Gestión", Pt(40), WHITE, True, FH), P("de Eventos", Pt(40), WHITE, True, FH)], ls=1.0)
tb(s, Inches(0.7), Inches(3.95), Inches(7.0), Inches(0.9),
   [P("Una solución de consola para registrar eventos, controlar cupos, vender entradas y consultar estadísticas.", Pt(14), SUBT)], ls=1.1)
for i, nom in enumerate(INTEGRANTES):
    cx = Inches(0.7 + i*1.95)
    rect(s, cx, Inches(5.25), Inches(1.8), Inches(0.42), EMER_D, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, cx, Inches(5.25), Inches(1.8), Inches(0.42), [P(nom, Pt(9.5), WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tb(s, Inches(0.7), Inches(6.7), Inches(11), Inches(0.3),
   [P("Paradigmas de la Programación · FP-UNA · 2026", Pt(10.5), FAINT)])

# ============ SLIDE 2 — El problema ============
s = slide(); bg(s, BG); header(s, "Contexto", "El problema que resolvemos", 2)
card(s, Inches(0.6), Inches(1.65), Inches(3.9), Inches(2.55), AMBER, "Riesgo operativo",
     "Gestionar eventos de forma manual aumenta el riesgo de datos incompletos, sobreventa de entradas y pérdida de control del estado.")
card(s, Inches(4.7), Inches(1.65), Inches(3.9), Inches(2.55), BLUE, "Qué se necesita",
     "Centralizar eventos, cupos, confirmaciones y métricas en un flujo simple, con validaciones para fechas y cantidades.")
card(s, Inches(8.8), Inches(1.65), Inches(3.9), Inches(2.55), EMER, "Resultado esperado",
     "Una herramienta clara y demostrable en consola, que da visibilidad inmediata del estado de la agenda.")
pic_fit(s, os.path.join(MK, "mockup_alta.png"), Inches(0.6), Inches(4.45), Inches(6.6), Inches(2.4))
card(s, Inches(7.4), Inches(4.45), Inches(5.3), Inches(2.4), EMER_D, "Mensaje clave",
     "El sistema no solo guarda eventos: aplica las reglas por sí mismo y permite decidir con datos confiables.")

# ============ SLIDE 3 — Nuestra solución ============
s = slide(); bg(s, BG); header(s, "Propuesta", "Nuestra solución: una agenda de eventos", 3)
rect(s, Inches(0.6), Inches(1.65), Inches(6.0), Inches(5.0), CARD, line=LINE, lw=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
pic_fit(s, os.path.join(MK, "mockup_menu.png"), Inches(0.75), Inches(1.8), Inches(5.7), Inches(4.7))
card(s, Inches(6.9), Inches(1.65), Inches(5.8), Inches(2.6), EMER, "Qué permite hacer",
     ["• Crear una agenda especializada por deporte",
      "• Cargar eventos con fecha, lugar, categoría y cupo",
      "• Confirmar eventos y vender entradas con control de cupo",
      "• Consultar estadísticas de ocupación y estado"])
card(s, Inches(6.9), Inches(4.45), Inches(5.8), Inches(2.2), BLUE, "Por qué es útil",
     "Permite ver al instante qué eventos están cargados, cuántas entradas se vendieron y qué capacidad queda disponible.")

# ============ SLIDE 4 — Funcionalidades ============
s = slide(); bg(s, BG); header(s, "Capacidades", "Todo lo que podés hacer", 4)
feats = ["Registrar eventos","Mostrar y consultar","Buscar por nombre","Filtrar por categoría",
         "Vender entradas","Confirmar eventos","Ver estadísticas","Análisis funcional"]
cols=4; cw=Inches(2.92); ch=Inches(1.65); gx=Inches(0.6); gy=Inches(1.75)
for i,f in enumerate(feats):
    r,c=divmod(i,cols); x=gx+c*(cw+Inches(0.13)); y=gy+r*(ch+Inches(0.22))
    rect(s,x,y,cw,ch,CARD,line=LINE,lw=Pt(0.75),shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    col=[EMER,BLUE2,AMBER,EMER_D][i%4]
    ch_o=rect(s,x+Inches(0.28),y+Inches(0.28),Inches(0.5),Inches(0.5),col,shape=MSO_SHAPE.OVAL)
    cf=ch_o.text_frame; cf.word_wrap=False; pp=cf.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER
    rr=pp.add_run(); rr.text=str(i+1); rr.font.size=Pt(15); rr.font.bold=True; rr.font.color.rgb=WHITE; rr.font.name=FH
    cf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tb(s,x+Inches(0.28),y+Inches(0.95),cw-Inches(0.5),Inches(0.55),[P(f,Pt(13),INK,True,FH)])

# ============ SLIDE 5 — Demostración funcional ============
s = slide(); bg(s, BG); header(s, "Demostración", "Recorrido del sistema", 5)
steps=[("1. Alta inicial","Creamos la agenda y cargamos eventos con sus datos. La fecha se valida antes de registrar.",EMER),
       ("2. Control de cupos","Vendemos entradas y el sistema impide superar la capacidad disponible.",BLUE),
       ("3. Seguimiento","Confirmamos eventos y consultamos estadísticas: ocupación, estados y categorías.",AMBER)]
for i,(t,b,col) in enumerate(steps):
    x=Inches(0.6+i*4.07)
    card(s,x,Inches(1.65),Inches(3.85),Inches(2.1),col,t,b)
pic_fit(s, os.path.join(MK, "mockup_venta.png"), Inches(0.6), Inches(4.0), Inches(7.0), Inches(2.8))
card(s, Inches(7.8), Inches(4.0), Inches(4.9), Inches(2.8), EMER_D, "Qué se ve en la demo",
     ["• Rechazo de una fecha inválida","• Venta aceptada dentro del cupo",
      "• Venta excedida rechazada","• Estadísticas y reglamento del deporte"])

# ============ SLIDE 6 — Capturas del sistema (galería) ============
s = slide(); bg(s, BG); header(s, "El sistema en acción", "Capturas del sistema", 6)
shots=[("Menú principal", "mockup_menu.png"),
       ("Carga y validación", "mockup_alta.png"),
       ("Estadísticas", "mockup_stats.png"),
       ("Módulo funcional", "mockup_funcional.png")]
gx=Inches(0.6); gy=Inches(1.7); cw=Inches(6.0); chh=Inches(2.55)
for i,(cap,img) in enumerate(shots):
    r,c=divmod(i,2); x=gx+c*(cw+Inches(0.13)); y=gy+r*(chh+Inches(0.18))
    rect(s,x,y,cw,chh,CARD,line=LINE,lw=Pt(0.75),shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s,x,y,cw,Inches(0.34),DARK)
    tb(s,x+Inches(0.2),y,cw-Inches(0.3),Inches(0.34),[P(cap,Pt(10.5),WHITE,True)],anchor=MSO_ANCHOR.MIDDLE)
    pic_fit(s,os.path.join(MK,img),x+Inches(0.12),y+Inches(0.42),cw-Inches(0.24),chh-Inches(0.54))

# ============ SLIDE 7 — Beneficios ============
s = slide(); bg(s, BG); header(s, "Valor", "Beneficios del sistema", 7)
bens=[("Control","El cupo, las ventas y los lugares disponibles se gestionan de forma consistente para evitar sobreventa.",EMER),
      ("Trazabilidad","Cada evento se identifica por nombre, fecha, lugar, categoría, estado y ventas realizadas.",BLUE),
      ("Visibilidad","Las estadísticas resumen total, confirmados, pendientes, ventas, ocupación y categorías.",AMBER),
      ("Calidad de datos","La fecha se valida en formato AAAA-MM-DD y los cupos deben ser positivos. Lo inválido se rechaza.",EMER_D),
      ("Simplicidad","La consola guía con un menú directo, pensado para demostrar el flujo sin herramientas externas.",BLUE2),
      ("Extensión","La lógica está separada de la interfaz: se puede sumar otra interfaz o deportes sin reescribir el motor.",EMER)]
for i,(t,b,col) in enumerate(bens):
    r,c=divmod(i,3); x=Inches(0.6+c*4.07); y=Inches(1.7+r*2.5)
    card(s,x,y,Inches(3.85),Inches(2.3),col,t,b,tsz=13.5,bsz=10.2)

# ============ SLIDE 8 — Evidencia ============
s = slide(); bg(s, BG); header(s, "Evidencia", "Resultados de la demostración", 8)
card(s, Inches(0.6), Inches(1.7), Inches(3.85), Inches(3.0), EMER, "Datos de la demo",
     ["Agenda: Liga Verano","Especialización: futsal","Evento: Final del Torneo","Cupo: 200 entradas","Venta: 150 entradas"], bsz=11)
card(s, Inches(4.7), Inches(1.7), Inches(3.85), Inches(3.0), BLUE, "Respuesta del sistema",
     ["• Rechaza fecha inválida","• Acepta venta dentro del cupo","• Rechaza venta excedida","• Calcula la ocupación","• Muestra el reglamento"], bsz=11)
card(s, Inches(8.8), Inches(1.7), Inches(3.9), Inches(3.0), AMBER, "Indicadores",
     ["Total de eventos: 2","Confirmados: 1 | Pendientes: 1","Vendidas: 150 / 280","Ocupación: 54%"], bsz=11)
pic_fit(s, os.path.join(MK, "mockup_stats.png"), Inches(0.6), Inches(4.9), Inches(7.0), Inches(1.9))
card(s, Inches(7.8), Inches(4.9), Inches(4.9), Inches(1.9), EMER_D, "Mensaje clave",
     "No depende de supuestos: la consola muestra la validación y los indicadores en vivo.")

# ============ SLIDE 9 — Cómo está construido + crecer ============
s = slide(); bg(s, BG); header(s, "Diseño", "Cómo está construido y hacia dónde crece", 9)
arch=[("Evento","Cada evento individual, con sus datos y sus reglas (cupo, ventas, estado).",EMER),
      ("Agenda de eventos","El contenedor que gestiona la colección: alta, búsqueda, filtrado y estadísticas.",BLUE),
      ("Agenda deportiva","Una agenda especializada en un deporte; reutiliza la base y agrega su reglamento.",AMBER)]
for i,(t,b,col) in enumerate(arch):
    card(s,Inches(0.6+i*4.07),Inches(1.7),Inches(3.85),Inches(2.0),col,t,b)
rect(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(2.7), DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.6), Inches(4.0), Inches(0.12), Inches(2.7), EMER)
tb(s, Inches(1.0), Inches(4.3), Inches(11.4), Inches(0.4), [P("Pensado para crecer", Pt(16), EMER, True, FH)])
grow=["La lógica está desacoplada de la interfaz: se podría sumar una interfaz gráfica o web reutilizando el mismo motor.",
      "Agregar un deporte nuevo o un tipo de agenda no requiere reescribir lo existente.",
      "Código validado, documentado y con diagrama técnico que respalda el diseño."]
y=4.85
for g in grow:
    rect(s, Inches(1.0), Inches(y+0.06), Inches(0.13), Inches(0.13), EMER, shape=MSO_SHAPE.OVAL)
    tb(s, Inches(1.3), Inches(y-0.06), Inches(11.1), Inches(0.5), [P(g, Pt(11.5), SUBT)])
    y+=0.58

# ============ SLIDE 10 — Cierre ============
s = slide(); bg(s, BG); header(s, "Cierre", "Entregable y próximos pasos", 10)
card(s, Inches(0.6), Inches(1.7), Inches(3.85), Inches(2.3), EMER, "Entregable actual",
     "Sistema funcional en consola para gestionar eventos, validar datos, controlar cupos y consultar estadísticas.")
card(s, Inches(4.7), Inches(1.7), Inches(3.85), Inches(2.3), BLUE, "Valor confirmado",
     "La demo comprueba el alta de eventos, el rechazo de datos inválidos, la venta controlada y el resumen de indicadores.")
card(s, Inches(8.8), Inches(1.7), Inches(3.9), Inches(2.3), AMBER, "Próximas mejoras",
     "Persistencia de datos, reportes exportables o una interfaz gráfica reutilizando el mismo motor.")
pic_fit(s, os.path.join(MK, "mockup_filtrar.png"), Inches(0.6), Inches(4.25), Inches(7.0), Inches(2.5))
card(s, Inches(7.8), Inches(4.25), Inches(4.9), Inches(2.5), EMER_D, "En síntesis",
     "El sistema está listo y demuestra el flujo completo; su diseño desacoplado permite que evolucione sin reescribir la lógica.")

# ============ SLIDE 11 — Gracias ============
s = slide(); bg(s, DARK)
rect(s, Inches(-1.0), Inches(4.5), Inches(4.5), Inches(4.5), RGBColor(0x12,0x1C,0x30), shape=MSO_SHAPE.OVAL)
rect(s, Inches(0.7), Inches(2.5), Inches(1.6), Inches(0.14), EMER)
tb(s, Inches(0.66), Inches(2.8), Inches(11), Inches(1.1), [P("¡Gracias!", Pt(46), WHITE, True, FH)])
tb(s, Inches(0.7), Inches(4.1), Inches(10.5), Inches(0.5),
   [P("Un sistema que centraliza, valida y analiza tus eventos.", Pt(16), SUBT)])
rect(s, Inches(0.72), Inches(4.95), Inches(7.4), Pt(1.2), RGBColor(0x2B,0x47,0x40))
tb(s, Inches(0.7), Inches(5.15), Inches(11), Inches(0.5), [P("¿Preguntas?", Pt(20), EMER, True, FH)])
tb(s, Inches(0.7), Inches(6.5), Inches(11), Inches(0.35), [P(PIE_CIERRE, Pt(10.5), FAINT)])

prs.save(OUT)
print("Generado:", OUT, "| slides:", len(prs.slides._sldIdLst))
